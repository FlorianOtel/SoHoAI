"""
Document parsing, chunking, and Qdrant upsert for the RAG pipeline.

Implements the 7-step atomic worker loop from RAG-strategy.md §4.4.
Called by the ingestion daemon (utils/rag_ingest_daemon.py) and the
FastAPI ingest endpoints in main.py.

Chunking strategy (RAG-strategy.md §3.4):
  - Parent-child (PDF, DOCX, IPYNB, long MD): child ~250 tok embedded;
    parent ~1000 tok stored as context returned to the LLM.
  - Flat 512-tok (PPTX, YAML, CSV, short docs): parent == child.

Parsing:
  - Structured formats (PDF, DOCX, PPTX): docling → markdown export.
  - PPTX fallback: OLE2 detection → LibreOffice conversion → python-pptx when docling fails.
  - IPYNB: dedicated cell extractor (JSON parse → markdown + code cells).
  - Text formats (TXT, MD, YAML, CSV): direct UTF-8 read.
  - docling failures fall back to python-pptx (PPTX) or raw UTF-8 read (others).
"""

from __future__ import annotations

import asyncio
import functools
import json
import logging
import re
import shutil
import subprocess
import tempfile
import threading
import urllib.request
import uuid
from datetime import datetime
from pathlib import Path

import tiktoken
from docling.document_converter import DocumentConverter
from qdrant_client import QdrantClient
from qdrant_client.models import (
    FieldCondition,
    Filter,
    FilterSelector,
    MatchValue,
    PointStruct,
    SparseVector,
)

from .collection import DOCUMENTS_COLLECTION, ensure_collection, collection_has_sparse
from .embeddings import embed_batch, embed_sparse
from .schema import (
    FIELD_CHUNK_INDEX,
    FIELD_FILE_NAME,
    FIELD_FILE_TYPE,
    FIELD_OWNER,
    FIELD_PAGE,
    FIELD_PARENT_TEXT,
    FIELD_PROJECT,
    FIELD_SESSION_ID,
    FIELD_SESSION_TITLE,
    FIELD_SOURCE_PATH,
    FIELD_TAG,
    FIELD_TEXT,
    SPARSE_VECTOR_NAME,
)
from .state import StateDB

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Sparse vector readiness cache
# ---------------------------------------------------------------------------

# Module-level cache: stores whether the Qdrant collection has sparse_text configured.
# Lazily initialized on first call to _is_sparse_ready() and cached for the lifetime
# of the worker process (across multiple file ingestions).
_sparse_ready: bool | None = None


def _is_sparse_ready(qdrant_client: QdrantClient) -> bool:
    """
    Check if the Qdrant collection has sparse vector support enabled.

    Lazy-initializes the _sparse_ready cache on first call. Returns True if the
    collection has the SPARSE_VECTOR_NAME (sparse_text) vector space configured;
    False before migration via rag_sparse_migrate.py or if the collection is absent.

    Args:
        qdrant_client: Qdrant client instance.

    Returns:
        True if sparse vectors are supported, False otherwise.
    """
    global _sparse_ready
    if _sparse_ready is None:
        _sparse_ready = collection_has_sparse(qdrant_client)
        logger.debug("Sparse vector readiness: %s", _sparse_ready)
    return _sparse_ready


# ---------------------------------------------------------------------------
# Chunking parameters (RAG-strategy.md §3.4)
# ---------------------------------------------------------------------------

_PARENT_CHUNK_SIZE    = 1000   # tokens — midpoint of 800–1200 range
_PARENT_CHUNK_OVERLAP = 100
_CHILD_CHUNK_SIZE     = 250    # tokens — midpoint of 200–300 range; safe under bge-m3's 8192-token context
_CHILD_CHUNK_OVERLAP  = 20
_FLAT_CHUNK_SIZE      = 512
_FLAT_CHUNK_OVERLAP   = 50

# Qdrant rejects upsert requests whose JSON payload exceeds its HTTP body limit.
# Large files (textbooks, big CSVs) can produce thousands of points whose combined
# parent_text fields push well past that limit. Splitting into smaller batches keeps
# each request well under the threshold while preserving all-or-nothing semantics
# at the file level (step 0 already deleted stale points before we get here).
_UPSERT_BATCH_SIZE = 256   # points per qdrant_client.upsert() call

# Always use flat chunking for these types (compact; parent-child overhead not worth it)
_ALWAYS_FLAT_TYPES = {"pptx", "yaml", "yml", "csv"}

# If total token count is below this threshold, use flat even for PDF/IPYNB/MD
_SHORT_DOC_TOKENS = 600

# Structured formats that go through docling; everything else is direct text read
_DOCLING_TYPES = {"pdf", "docx", "pptx"}   # ipynb handled by _parse_ipynb()

# tiktoken cl100k_base — consistent encoding across ingest and search
_enc = tiktoken.get_encoding("cl100k_base")

# Thread-local DocumentConverter — one instance per OS thread so concurrent
# asyncio.to_thread() calls don't share docling's internal conversion state.
_thread_local = threading.local()


def _get_converter() -> DocumentConverter:
    if not hasattr(_thread_local, "converter"):
        _thread_local.converter = DocumentConverter()
    return _thread_local.converter


# ---------------------------------------------------------------------------
# Parsing
# ---------------------------------------------------------------------------

_UUID_RE = re.compile(
    r"^[0-9a-f]{8}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{12}$"
)


@functools.lru_cache(maxsize=8)
def _build_title_map(dot_claude_dir: Path) -> dict[str, str]:
    """Walk ~/.claude/chats/ symlinks → {realpath: human_title}.
    Excludes UUID-named symlinks (auto-created by Claude Code, not user-chosen).
    Returns {} if chats dir absent or inaccessible.
    """
    chats_dir = dot_claude_dir / "chats"
    if not chats_dir.is_dir():
        return {}
    title_map: dict[str, str] = {}
    try:
        for category in chats_dir.iterdir():
            if not category.is_dir():
                continue
            for f in category.iterdir():
                if f.suffix != ".jsonl" or not f.is_symlink():
                    continue
                if _UUID_RE.match(f.stem):
                    continue
                try:
                    title_map[str(f.resolve())] = f.stem
                except OSError:
                    continue
    except OSError as exc:
        logger.warning("_build_title_map: cannot scan %s: %s", chats_dir, exc)
    return title_map


def _parse_ipynb(path: Path) -> str:
    """
    Extract clean text from a Jupyter notebook.

    Concatenates markdown and code cells in order, separated by blank lines.
    Skips raw cells and empty cells. Cell outputs are ignored — source only.
    Falls back to raw UTF-8 read if the file is not valid JSON.
    """
    try:
        nb = json.loads(path.read_text(encoding="utf-8", errors="replace"))
    except json.JSONDecodeError:
        return path.read_text(encoding="utf-8", errors="replace")

    parts: list[str] = []
    for cell in nb.get("cells", []):
        cell_type = cell.get("cell_type", "")
        source = cell.get("source", [])
        text = "".join(source).strip() if isinstance(source, list) else str(source).strip()
        if not text:
            continue
        if cell_type == "markdown":
            parts.append(text)
        elif cell_type == "code":
            parts.append(f"```python\n{text}\n```")
    return "\n\n".join(parts)


# OLE2 magic bytes — Composite Document File V2 (binary .ppt format)
_OLE2_MAGIC = b"\xD0\xCF\x11\xE0"


def _synthesize_title(path: Path, max_chars: int = 60) -> str:
    """Derive a title from the first non-empty, non-command user message.
    Falls back to path.stem (UUID) on any failure.
    """
    try:
        lines = path.read_text(encoding="utf-8", errors="replace").splitlines()
    except OSError:
        return path.stem
    for raw in lines:
        try:
            entry = json.loads(raw.strip())
        except (json.JSONDecodeError, ValueError):
            continue
        if entry.get("type") != "user" or entry.get("isMeta"):
            continue
        content = (entry.get("message") or {}).get("content", "")
        if isinstance(content, str):
            text = content.strip()
        elif isinstance(content, list):
            text = " ".join(
                b.get("text", "").strip()
                for b in content
                if isinstance(b, dict) and b.get("type") == "text"
            )
        else:
            text = ""
        if not text or text.startswith("<"):
            continue
        title = text.replace("\n", " ").strip()[:max_chars]
        if len(text) > max_chars:
            title += "..."
        return title
    return path.stem


def _parse_claude_chat(path: Path) -> tuple[str, dict]:
    """
    Extract user/assistant text turns from a Claude Code session JSONL.

    Returns (text, metadata) where metadata contains session_id and project
    for storage as Qdrant payload fields.

    Skips: meta entries, tool_use blocks, tool_result blocks, thinking blocks,
    XML-tagged command strings, and file-history-snapshot entries.
    """
    import json as _json

    session_id = ""
    cwd = ""
    timestamp = ""
    model = ""
    git_branch = ""
    version = ""

    turns: list[str] = []

    try:
        lines = path.read_text(encoding="utf-8", errors="replace").splitlines()
    except OSError:
        return "", {}

    for raw_line in lines:
        raw_line = raw_line.strip()
        if not raw_line:
            continue
        try:
            entry = _json.loads(raw_line)
        except _json.JSONDecodeError:
            continue

        entry_type = entry.get("type", "")

        # Collect session metadata from any entry
        if not session_id:
            session_id = entry.get("sessionId", "")
        if not cwd:
            cwd = entry.get("cwd", "")
        if not git_branch:
            git_branch = entry.get("gitBranch", "")
        if not version:
            version = entry.get("version", "")

        if entry_type == "user":
            if entry.get("isMeta"):
                continue
            msg = entry.get("message", {})
            content = msg.get("content", "")
            text = ""
            if isinstance(content, str):
                c = content.strip()
                if c and not c.startswith("<"):
                    text = c
            elif isinstance(content, list):
                parts = []
                for block in content:
                    if isinstance(block, dict) and block.get("type") == "text":
                        t = block.get("text", "").strip()
                        if t:
                            parts.append(t)
                text = "\n".join(parts)
            if text:
                if not timestamp:
                    timestamp = entry.get("timestamp", "")
                turns.append(f"**User**: {text}")

        elif entry_type == "assistant":
            msg = entry.get("message", {})
            if not model:
                model = msg.get("model", "")
            content = msg.get("content", [])
            if isinstance(content, list):
                parts = []
                for block in content:
                    if isinstance(block, dict) and block.get("type") == "text":
                        t = block.get("text", "").strip()
                        if t:
                            parts.append(t)
                text = "\n".join(parts)
                if text:
                    turns.append(f"**Claude**: {text}")

    if not turns:
        return "", {}

    project = Path(cwd).name if cwd else "unknown"

    header = (
        f"# Claude Code Session\n\n"
        f"Session: {session_id}\n"
        f"Project: {project}  ({cwd})\n"
        f"Date: {timestamp}\n"
        f"Model: {model}\n"
        f"Branch: {git_branch}\n"
        f"Version: {version}\n\n"
        f"---\n\n"
    )

    body = "\n\n".join(turns)
    text = header + body

    # --- session_title derivation ---
    is_subagent = path.parent.name == "subagents"
    if is_subagent:
        # path: .../projects/<mangled>/<uuid>/subagents/agent-xxx.jsonl
        parent_jsonl = path.parent.parent.parent / (path.parent.parent.name + ".jsonl")
        dot_claude_dir = path.parents[4]
        title_map = _build_title_map(dot_claude_dir)
        parent_title = title_map.get(str(parent_jsonl.resolve())) or _synthesize_title(parent_jsonl)
        session_title = f"{parent_title} [subagent]"
    else:
        # path: .../projects/<mangled>/<uuid>.jsonl
        dot_claude_dir = path.parents[2]
        title_map = _build_title_map(dot_claude_dir)
        session_title = title_map.get(str(path.resolve())) or _synthesize_title(path)

    metadata = {
        FIELD_SESSION_ID:    session_id,
        FIELD_PROJECT:       project,
        FIELD_SESSION_TITLE: session_title,
    }
    return text, metadata


def _parse_opencode_session(path: str, oc_api_url: str) -> tuple[str, dict]:
    """
    Fetch an opencode session via HTTP and extract user/assistant text turns.

    path is the synthetic key "opencode://{session_id}".
    oc_api_url is the opencode API base URL (e.g. "http://localhost:4096").

    Returns (text, metadata) with session_id, project, session_title.
    On any HTTP failure, returns ("", {}) to signal parse failure.
    """
    # Parse session_id from synthetic path
    if not path.startswith("opencode://"):
        return "", {}
    session_id = path[len("opencode://"):]
    if not session_id:
        return "", {}

    turns: list[str] = []
    session_title = ""
    project = "unknown"
    created_timestamp = ""

    try:
        # Fetch session metadata (for title and created timestamp)
        session_url = f"{oc_api_url}/session/{session_id}"
        with urllib.request.urlopen(session_url, timeout=5) as response:
            session = json.load(response)
        session_title = session.get("title", "").strip()
        project = Path(session.get("directory", "")).name or "unknown"
        time_info = session.get("time", {})
        created_ms = time_info.get("created", 0)
        if created_ms:
            created_timestamp = datetime.fromtimestamp(created_ms / 1000.0).isoformat()

        # Fetch messages
        messages_url = f"{oc_api_url}/session/{session_id}/message"
        with urllib.request.urlopen(messages_url, timeout=5) as response:
            messages = json.load(response)
    except Exception:
        return "", {}

    # Extract text turns from messages
    first_user_text = ""
    for message in messages:
        info = message.get("info", {})
        role = info.get("role", "")
        parts = message.get("parts", [])

        # Extract text from parts; skip non-text types
        text_parts = []
        for part in parts:
            part_type = part.get("type", "")
            if part_type == "text":
                part_text = part.get("text", "").strip()
                if part_text:
                    text_parts.append(part_text)

        if not text_parts:
            continue

        combined_text = "\n".join(text_parts)

        # Format as user/assistant turn
        if role == "user":
            if not combined_text.startswith("/"):  # Skip slash commands
                turns.append(f"**User**: {combined_text}")
                if not first_user_text:
                    first_user_text = combined_text
        elif role == "assistant":
            turns.append(f"**Assistant**: {combined_text}")

    if not turns:
        return "", {}

    # Synthesise session_title if absent or empty
    if not session_title:
        if first_user_text:
            # Use first non-empty user text (up to 60 chars), stripping leading <
            text = first_user_text.lstrip("<").replace("\n", " ").strip()[:60]
            session_title = text + "..." if len(first_user_text) > 60 else text
        else:
            session_title = session_id  # Fallback to session_id

    # Build header block
    header = (
        f"# Opencode Session\n\n"
        f"Session: {session_id}\n"
        f"Project: {project}\n"
        f"Date: {created_timestamp}\n"
        f"Title: {session_title}\n\n"
        f"---\n\n"
    )

    body = "\n\n".join(turns)
    text = header + body

    metadata = {
        FIELD_SESSION_ID:    session_id,
        FIELD_PROJECT:       project,
        FIELD_SESSION_TITLE: session_title,
    }
    return text, metadata


def _parse_pptx(path: Path) -> str:
    """
    Extract text from a PowerPoint file using python-pptx.

    Used as fallback when docling fails to detect the PPTX format. Handles
    two cases:

    1. OLE2 binary .ppt disguised as .pptx (Composite Document File V2,
       detected by magic bytes): convert to real .pptx via LibreOffice
       headless, then parse with python-pptx.
    2. Genuine OpenXML .pptx that docling rejected: parse directly with
       python-pptx.

    Each slide is prefixed with "Slide N" so chunk boundaries remain
    interpretable during retrieval. Temp files are always cleaned up.
    """
    from pptx import Presentation  # docling dep — always present in this venv

    # Detect OLE2 binary .ppt by magic bytes
    with open(path, "rb") as fh:
        magic = fh.read(4)

    tmp_dir: str | None = None
    pptx_path = path
    if magic == _OLE2_MAGIC:
        logger.info(
            "OLE2 binary .ppt detected for %s — attempting LibreOffice conversion", path.name
        )
        tmp_dir = tempfile.mkdtemp()
        try:
            result = subprocess.run(
                [
                    "libreoffice", "--headless",
                    "--convert-to", "pptx",
                    "--outdir", tmp_dir,
                    str(path),
                ],
                check=False,
                timeout=120,
                capture_output=True,
            )
            pptx_path = Path(tmp_dir) / f"{path.stem}.pptx"
            if result.returncode != 0 or not pptx_path.exists():
                stderr = result.stderr.decode(errors="replace").strip()
                raise ValueError(
                    f"LibreOffice could not convert OLE2 .ppt '{path.name}' "
                    f"(file may be IRM/DRM-encrypted or corrupted): "
                    f"{stderr or 'no stderr'}"
                )
        except Exception:
            shutil.rmtree(tmp_dir, ignore_errors=True)
            raise

    try:
        prs = Presentation(str(pptx_path))
        parts: list[str] = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_lines: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        slide_lines.append(text)
            if slide_lines:
                parts.append(f"Slide {slide_num}\n" + "\n".join(slide_lines))
        return "\n\n".join(parts)
    finally:
        if tmp_dir:
            shutil.rmtree(tmp_dir, ignore_errors=True)


def _parse_to_text(file_path: str, file_type: str) -> str:
    """
    Extract full text from a document.

    Structured formats (pdf, docx, pptx) go through docling → Markdown.
    PPTX fallback: python-pptx when docling format detection fails.
    IPYNB: dedicated cell extractor (avoids docling which doesn't support it).
    Text formats (txt, md, yaml, csv): direct UTF-8 read.

    Synchronous and potentially slow for large PDFs; callers use asyncio.to_thread().
    """
    path = Path(file_path)
    if file_type == "ipynb":
        return _parse_ipynb(path)
    if file_type == "jsonl":
        text, _ = _parse_claude_chat(path)
        return text
    if file_type not in _DOCLING_TYPES:
        return path.read_text(encoding="utf-8", errors="replace")

    try:
        result = _get_converter().convert(source=str(path))
        return result.document.export_to_markdown()
    except Exception as exc:
        if file_type == "pptx":
            logger.warning(
                "docling failed for %s (%s) — trying python-pptx fallback: %s",
                path.name, file_type, exc,
            )
            # No raw-read fallback for binary PPTX: if _parse_pptx raises, let it
            # propagate so ingest_file() marks the file failed rather than embedding garbage.
            return _parse_pptx(path)
        logger.warning(
            "docling failed for %s (%s) — falling back to raw read: %s",
            path.name, file_type, exc,
        )
        return path.read_text(encoding="utf-8", errors="replace")


# ---------------------------------------------------------------------------
# Token-based chunking
# ---------------------------------------------------------------------------

def _split_tokens(text: str, chunk_size: int, overlap: int) -> list[str]:
    """Split text into token-bounded chunks with overlap. Preserves order."""
    tokens = _enc.encode(text)
    if not tokens:
        return []
    chunks: list[str] = []
    start = 0
    while start < len(tokens):
        end = min(start + chunk_size, len(tokens))
        chunks.append(_enc.decode(tokens[start:end]))
        if end == len(tokens):
            break
        start += chunk_size - overlap
    return chunks


def _chunk_parent_child(text: str) -> list[tuple[str, str]]:
    """
    Split text using the parent-child strategy.

    Each parent (~1000 tok) is split into smaller children (~250 tok).
    The child is what gets embedded (precise similarity); the parent is
    stored in the Qdrant payload and returned to the LLM for richer context.

    Returns: list of (child_text, parent_text) pairs.
    """
    pairs: list[tuple[str, str]] = []
    for parent in _split_tokens(text, _PARENT_CHUNK_SIZE, _PARENT_CHUNK_OVERLAP):
        for child in _split_tokens(parent, _CHILD_CHUNK_SIZE, _CHILD_CHUNK_OVERLAP):
            pairs.append((child, parent))
    return pairs


def _chunk_flat(text: str) -> list[tuple[str, str]]:
    """
    Flat chunking for compact files (PPTX, short TXT, YAML, CSV).

    Returns (chunk, chunk) pairs — parent == child (no extra context layer).
    """
    return [(c, c) for c in _split_tokens(text, _FLAT_CHUNK_SIZE, _FLAT_CHUNK_OVERLAP)]


def _select_chunks(text: str, file_type: str) -> tuple[list[tuple[str, str]], str]:
    """
    Choose chunking strategy and return (pairs, strategy_label).

    Flat when:
      - file_type is in _ALWAYS_FLAT_TYPES (pptx, yaml, csv), OR
      - document is short (< _SHORT_DOC_TOKENS tokens).
    Parent-child for everything else.
    """
    if file_type in _ALWAYS_FLAT_TYPES:
        return _chunk_flat(text), "flat"
    if len(_enc.encode(text)) < _SHORT_DOC_TOKENS:
        return _chunk_flat(text), "flat(short)"
    return _chunk_parent_child(text), "parent-child"


# ---------------------------------------------------------------------------
# Worker loop
# ---------------------------------------------------------------------------

async def ingest_file(
    file_path: str,
    owner: str,
    rag_cfg: dict,
    state_db: StateDB,
    qdrant_client: QdrantClient,
    tag: str = "",
    embed_concurrency: int = 5,
) -> None:
    """
    Atomically parse, chunk, embed, and upsert one document into Qdrant.

    Implements the 7-step atomic worker loop from RAG-strategy.md §4.4:
      0. Delete stale Qdrant points for this source_path (idempotency).
      1. Mark file as 'processing' in StateDB.
      2. Parse document → full text.
      3. Chunk text → (child, parent) pairs.
      4. Embed child texts via Ollama.
      5. Build PointStruct list with full payload.
      6. Single atomic upsert to Qdrant.
      7. Mark file as 'completed' in StateDB.

    On any failure: marks the file as failed (with auto-retry logic in
    StateDB) and re-raises so the daemon loop can move to the next file.

    Args:
        file_path:        Absolute NFS path to the file.
        owner:            Owner string (e.g. "florian") — from StateDB row.
        rag_cfg:          config["rag"] dict — embedding_model, ollama_url.
        state_db:         StateDB instance for ingestion status tracking.
        qdrant_client:    Qdrant client in local persistent mode (NAS).
        tag:              Optional tag derived by the NFS scanner (e.g. "certifications").
        embed_concurrency: Max parallel Ollama embedding requests per file.
                          Controlled by --batch in rag_ingest_daemon.py (default 5).
    """
    model      = rag_cfg.get("embedding_model", "bge-m3")
    ollama_url = rag_cfg.get("ollama_url", "http://192.168.1.93:11434/api/embeddings")

    file_name     = Path(file_path).name
    raw_file_type = Path(file_path).suffix.lstrip(".").lower()

    try:
        # ------------------------------------------------------------------
        # Step 0: delete stale Qdrant points (delete-before-insert)
        #
        # Ensures re-ingestion of a modified file never leaves duplicate points.
        # Even if the collection is brand-new, deleting 0 matching points is valid.
        # ------------------------------------------------------------------
        ensure_collection(qdrant_client)
        qdrant_client.delete(
            collection_name=DOCUMENTS_COLLECTION,
            points_selector=FilterSelector(
                filter=Filter(
                    must=[
                        FieldCondition(
                            key=FIELD_SOURCE_PATH,
                            match=MatchValue(value=file_path),
                        )
                    ]
                )
            ),
        )

        # ------------------------------------------------------------------
        # Step 1: mark processing
        # ------------------------------------------------------------------
        state_db.mark_processing(file_path)

        # ------------------------------------------------------------------
        # Step 2: parse document
        #
        # Claude Code session files (.jsonl) are parsed directly to extract
        # user/assistant turns and cross-reference metadata (session_id, project).
        # All other formats go through _parse_to_text (docling / ipynb / utf-8).
        # asyncio.to_thread keeps the event loop responsive for other tasks.
        # ------------------------------------------------------------------
        state_db.set_progress(file_path, "parsing")
        if file_path.startswith("opencode://"):
            oc_api_url = rag_cfg.get("opencode_api_url", "http://localhost:4096")
            text, chat_meta = await asyncio.to_thread(_parse_opencode_session, file_path, oc_api_url)
            file_type = "opencode"
        elif raw_file_type == "jsonl":
            text, chat_meta = await asyncio.to_thread(_parse_claude_chat, Path(file_path))
            file_type = "claude_chat"
        else:
            text = await asyncio.to_thread(_parse_to_text, file_path, raw_file_type)
            file_type = raw_file_type
            chat_meta = {}

        # For chat sessions, the raw file_name is a meaningless identifier
        # (UUID-named .jsonl for Claude Code, ses_<id> for opencode). Replace it
        # with the human-readable session_title so any renderer using file_name
        # gets a sensible label without per-type special-casing.
        if chat_meta.get(FIELD_SESSION_TITLE):
            file_name = chat_meta[FIELD_SESSION_TITLE]

        if not text.strip():
            logger.warning("Empty document, skipping: %s", file_path)
            state_db.mark_completed(file_path)
            return

        # Skip documents whose parsed text is mostly binary/garbage characters.
        # Corrupt scanned PDFs and binary-encoded files produce noise vectors
        # that crowd out legitimate search results.
        _printable = sum(1 for c in text if c.isprintable()) / max(len(text), 1)
        if _printable < 0.85:
            logger.warning(
                "Skipping %s — parsed text is %.0f%% printable (binary/corrupt content)",
                file_path, _printable * 100,
            )
            state_db.mark_ignored(file_path, f"binary content: {_printable:.0%} printable chars")
            return

        # ------------------------------------------------------------------
        # Step 3: chunk
        # ------------------------------------------------------------------
        pairs, strategy = _select_chunks(text, file_type)
        n = len(pairs)
        state_db.set_progress(file_path, f"chunking ({n} chunks)")
        logger.info(
            "Chunked %s → %d chunks  strategy=%s  file_type=%s",
            file_name, n, strategy, file_type,
        )

        if not pairs:
            logger.warning("No chunks produced, skipping: %s", file_path)
            state_db.mark_completed(file_path)
            return

        # ------------------------------------------------------------------
        # Step 4: embed child texts via Ollama (embed_concurrency-concurrent semaphore)
        # ------------------------------------------------------------------
        child_texts = [child for child, _ in pairs]

        def _log_embed_progress(done: int, total: int) -> None:
            logger.info("Embedding progress: %d/%d  %s", done, total, file_name)

        state_db.set_progress(file_path, f"embedding 0/{n}")
        vectors = await embed_batch(
            child_texts, model=model, ollama_url=ollama_url,
            progress_cb=_log_embed_progress,
            concurrency=embed_concurrency,
        )
        state_db.set_progress(file_path, f"embedding {n}/{n}")

        # ------------------------------------------------------------------
        # Step 5: build PointStructs
        #
        # Random UUIDs — step 0 handles cleanup before re-insert, so there
        # are no orphan-ID concerns when chunking produces a different count.
        #
        # Sparse vectors (BM25) are conditionally included if the collection
        # has sparse_text configured. Before migration via rag_sparse_migrate.py,
        # ingest only dense vectors (backward compatible).
        # ------------------------------------------------------------------
        sparse_ready = _is_sparse_ready(qdrant_client)
        points = []
        for chunk_idx, ((child_text, parent_text), vector) in enumerate(zip(pairs, vectors)):
            if sparse_ready:
                sparse_indices, sparse_values = embed_sparse(child_text)
                point_vector = {
                    "": vector,  # "" is the Qdrant key for the unnamed dense vector
                    SPARSE_VECTOR_NAME: SparseVector(indices=sparse_indices, values=sparse_values),
                }
            else:
                point_vector = vector  # dense-only until migration is complete

            points.append(
                PointStruct(
                    id=str(uuid.uuid4()),
                    vector=point_vector,
                    payload={
                        FIELD_TEXT:        child_text,
                        FIELD_PARENT_TEXT: parent_text,
                        FIELD_OWNER:       owner,
                        FIELD_SOURCE_PATH: file_path,
                        FIELD_FILE_NAME:   file_name,
                        FIELD_FILE_TYPE:   file_type,
                        FIELD_PAGE:        0,   # page-level tracking: future enhancement
                        FIELD_CHUNK_INDEX: chunk_idx,
                        FIELD_TAG:         tag,
                        **chat_meta,
                    },
                )
            )

        # ------------------------------------------------------------------
        # Step 6: batched upsert
        #
        # A single upsert with all points can exceed Qdrant's HTTP body limit
        # for large files (Qdrant returns HTTP 400 "Payload error: JSON payload
        # (N bytes)"). Splitting into _UPSERT_BATCH_SIZE-point batches keeps
        # each request well under the limit.
        #
        # Atomicity note: step 0 deleted all stale points for this file before
        # processing began, so the collection is already clean. Individual batch
        # failures leave a partial set of points for this file — the daemon's
        # mark_failed() + retry logic will re-run from step 0 on the next
        # attempt, deleting the partial set before re-inserting cleanly.
        # ------------------------------------------------------------------
        state_db.set_progress(file_path, "upserting")
        for batch_start in range(0, len(points), _UPSERT_BATCH_SIZE):
            qdrant_client.upsert(
                collection_name=DOCUMENTS_COLLECTION,
                points=points[batch_start : batch_start + _UPSERT_BATCH_SIZE],
            )

        # ------------------------------------------------------------------
        # Step 7: mark completed
        # ------------------------------------------------------------------
        state_db.mark_completed(file_path)
        logger.info(
            "Ingested %s: %d points  owner=%s  tag=%s",
            file_name, len(points), owner, tag or "—",
        )

    except Exception as exc:
        logger.error("Ingestion failed for %s: %s", file_path, exc, exc_info=True)
        state_db.mark_failed(file_path, str(exc))
        raise
